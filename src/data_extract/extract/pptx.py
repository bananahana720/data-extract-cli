"""PPTX extractor adapter."""

from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, Tuple

from data_extract.extract.adapter import ExtractorAdapter


class PptxExtractorAdapter(ExtractorAdapter):
    """Extractor for Microsoft PowerPoint presentations."""

    def __init__(self) -> None:
        super().__init__(format_name="pptx")

    def extract(self, file_path: Path) -> Tuple[str, Dict[str, Any], Dict[str, float]]:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise RuntimeError("python-pptx is required for PPTX extraction") from exc

        presentation = Presentation(str(file_path))
        lines = []
        text_box_count = 0
        notes_count = 0

        for idx, slide in enumerate(presentation.slides, start=1):
            lines.append(f"# Slide {idx}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    content = shape.text.strip()
                    if content:
                        lines.append(content)
                        text_box_count += 1

            notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
            notes = notes.strip() if notes else ""
            if notes:
                lines.append(f"Notes: {notes}")
                notes_count += 1

        text = "\n\n".join(lines)
        structure = {
            "slide_count": len(presentation.slides),
            "text_box_count": text_box_count,
            "notes_count": notes_count,
        }
        quality = {
            "extraction_confidence": 1.0,
        }
        return text, structure, quality
